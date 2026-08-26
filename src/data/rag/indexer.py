"""Build and query a local RAG knowledge index.

The first implementation is dependency-light and deterministic. It uses a
lexical TF-IDF style sparse vector so the pipeline works on a new machine
before API keys or embedding services are configured. The JSON index format is
kept explicit so an OpenAI/vector-store backend can be added later without
changing the source collection layer.
"""

from __future__ import annotations

import json
import math
import re
import zipfile
from collections import Counter
from collections.abc import Mapping
from dataclasses import asdict, dataclass, field, replace
from datetime import datetime, timezone
from pathlib import Path
from typing import SupportsFloat, SupportsIndex, SupportsInt
from xml.etree import ElementTree

from src.data.rag.collector import chunk_text, extract_pdf_text, normalize_text

TOKEN_RE = re.compile(r"[A-Za-z][A-Za-z0-9_+\-]{1,}|[가-힣]{2,}|\d+(?:\.\d+)?")
PROJECT_ROOT = Path(__file__).resolve().parents[3]
DEFAULT_INDEX_PATH = PROJECT_ROOT / "data" / "rag" / "knowledge_index.json"
DEFAULT_ONLINE_CHUNKS_DIR = PROJECT_ROOT / "data" / "rag" / "online_corpus" / "chunks"


INTERNAL_SOURCE_PATTERNS: tuple[str, ...] = (
    "data/PPT/*.pptx",
    "data/datasets/DD/*.pdf",
    "data/datasets/DD/*.pptx",
    "data/datasets/DD/*.md",
    "docs/DD_Laminate*.md",
    "docs/*dd*.md",
    "docs/*Laminate*.md",
    "reports/dd_response*/**/*.md",
    "reports/dd_u3*/**/*.md",
    "research/papers/*.md",
    "research/syntheses/*.md",
    "research/recommendations/*.md",
)


@dataclass(frozen=True)
class KnowledgeChunk:
    """One searchable chunk in the local knowledge index."""

    chunk_id: str
    source_id: str
    title: str
    text: str
    source_kind: str
    source_path: str = ""
    url: str = ""
    topic: str = "composites"
    tags: list[str] = field(default_factory=list)
    chunk_index: int = 0
    token_counts: dict[str, float] = field(default_factory=dict)
    norm: float = 0.0

    def to_dict(self) -> dict[str, object]:
        return asdict(self)

    @classmethod
    def from_dict(cls, payload: dict[str, object]) -> KnowledgeChunk:
        return cls(
            chunk_id=str(payload["chunk_id"]),
            source_id=str(payload["source_id"]),
            title=str(payload["title"]),
            text=str(payload["text"]),
            source_kind=str(payload["source_kind"]),
            source_path=str(payload.get("source_path", "")),
            url=str(payload.get("url", "")),
            topic=str(payload.get("topic", "composites")),
            tags=_string_list_from_json(payload.get("tags", [])),
            chunk_index=_int_from_json(payload.get("chunk_index", 0)),
            token_counts=_float_mapping_from_json(payload.get("token_counts", {})),
            norm=_float_from_json(payload.get("norm", 0.0)),
        )


@dataclass(frozen=True)
class QueryResult:
    """One ranked retrieval result."""

    score: float
    chunk: KnowledgeChunk

    def to_dict(self) -> dict[str, object]:
        payload = self.chunk.to_dict()
        payload["score"] = self.score
        return payload


def build_knowledge_index(
    *,
    project_root: str | Path = PROJECT_ROOT,
    output_path: str | Path = DEFAULT_INDEX_PATH,
    include_online: bool = True,
    include_internal: bool = True,
) -> dict[str, object]:
    """Build a local searchable knowledge index from online and internal chunks."""
    root = Path(project_root)
    chunks: list[KnowledgeChunk] = []
    if include_online:
        chunks.extend(load_online_chunks(root / "data" / "rag" / "online_corpus" / "chunks"))
    if include_internal:
        chunks.extend(load_internal_chunks(root))

    chunks = attach_tfidf_vectors(chunks)
    payload = {
        "schema_version": "local-tfidf-v1",
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "chunk_count": len(chunks),
        "source_count": len({chunk.source_id for chunk in chunks}),
        "chunks": [chunk.to_dict() for chunk in chunks],
    }
    index_path = Path(output_path)
    index_path.parent.mkdir(parents=True, exist_ok=True)
    index_path.write_text(
        json.dumps(payload, indent=2, ensure_ascii=False) + "\n", encoding="utf-8"
    )
    return payload


def query_index(index_path: str | Path, query: str, *, top_k: int = 5) -> list[QueryResult]:
    """Query a local knowledge index."""
    index = json.loads(Path(index_path).read_text(encoding="utf-8"))
    chunks = [KnowledgeChunk.from_dict(item) for item in index.get("chunks", [])]
    query_counts = Counter(tokenize(query))
    if not query_counts:
        return []

    document_frequency = Counter[str]()
    for chunk in chunks:
        document_frequency.update(chunk.token_counts.keys())
    total_documents = max(len(chunks), 1)
    query_vector = {
        token: count * idf(token, document_frequency.get(token, 0), total_documents)
        for token, count in query_counts.items()
    }
    query_norm = vector_norm(query_vector)
    if query_norm == 0:
        return []

    results: list[QueryResult] = []
    for chunk in chunks:
        score = cosine_score(query_vector, query_norm, chunk.token_counts, chunk.norm)
        if score > 0:
            results.append(QueryResult(score=score, chunk=chunk))
    return sorted(results, key=lambda item: item.score, reverse=True)[:top_k]


def load_online_chunks(chunks_dir: str | Path) -> list[KnowledgeChunk]:
    root = Path(chunks_dir)
    if not root.exists():
        return []
    chunks: list[KnowledgeChunk] = []
    for path in sorted(root.glob("*.jsonl")):
        for line in path.read_text(encoding="utf-8").splitlines():
            if not line.strip():
                continue
            payload = json.loads(line)
            source_id = str(payload["source_id"])
            chunk_index = int(payload.get("chunk_index", len(chunks)))
            text = normalize_text(str(payload.get("text", "")))
            if not text:
                continue
            chunks.append(
                KnowledgeChunk(
                    chunk_id=f"online:{source_id}:{chunk_index}",
                    source_id=source_id,
                    title=str(payload.get("title", source_id)),
                    text=text,
                    source_kind="online",
                    url=str(payload.get("url", "")),
                    topic=str(payload.get("topic", "composites")),
                    tags=[str(tag) for tag in payload.get("tags", [])],
                    chunk_index=chunk_index,
                )
            )
    return chunks


def load_internal_chunks(project_root: str | Path) -> list[KnowledgeChunk]:
    root = Path(project_root)
    files: list[Path] = []
    for pattern in INTERNAL_SOURCE_PATTERNS:
        files.extend(root.glob(pattern))
    unique_files = sorted({path.resolve() for path in files if path.is_file()})

    chunks: list[KnowledgeChunk] = []
    for path in unique_files:
        text, note = extract_internal_text(path)
        if not text:
            continue
        source_id = internal_source_id(root, path)
        title = path.stem.replace("_", " ").replace("-", " ")
        tags = internal_tags(path)
        if note:
            tags.append(note)
        for index, chunk in enumerate(chunk_text(text)):
            chunks.append(
                KnowledgeChunk(
                    chunk_id=f"internal:{source_id}:{index}",
                    source_id=source_id,
                    title=title,
                    text=chunk,
                    source_kind="internal",
                    source_path=str(path),
                    topic=internal_topic(path),
                    tags=tags,
                    chunk_index=index,
                )
            )
    return chunks


def extract_internal_text(path: Path) -> tuple[str, str]:
    suffix = path.suffix.lower()
    if suffix in {".md", ".txt"}:
        return normalize_text(path.read_text(encoding="utf-8", errors="replace")), ""
    if suffix == ".pdf":
        return extract_pdf_text(path.read_bytes())
    if suffix == ".pptx":
        return extract_pptx_text(path), ""
    if suffix == ".json":
        return normalize_text(path.read_text(encoding="utf-8", errors="replace")), ""
    return "", f"unsupported:{suffix}"


def extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return extract_pptx_text_from_zip(path)

    presentation = Presentation(str(path))
    parts: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                slide_parts.append(text)
        if slide_parts:
            parts.append(f"Slide {slide_number}: {' '.join(slide_parts)}")
    return normalize_text("\n".join(parts))


def extract_pptx_text_from_zip(path: Path) -> str:
    namespaces = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    parts: list[str] = []
    with zipfile.ZipFile(path) as archive:
        slide_names = sorted(
            name for name in archive.namelist() if name.startswith("ppt/slides/slide")
        )
        for slide_number, name in enumerate(slide_names, start=1):
            root = ElementTree.fromstring(archive.read(name))
            texts = [node.text or "" for node in root.findall(".//a:t", namespaces)]
            if texts:
                parts.append(f"Slide {slide_number}: {' '.join(texts)}")
    return normalize_text("\n".join(parts))


def attach_tfidf_vectors(chunks: list[KnowledgeChunk]) -> list[KnowledgeChunk]:
    document_frequency = Counter[str]()
    raw_counts: list[Counter[str]] = []
    for chunk in chunks:
        counts = Counter(tokenize(chunk.text))
        raw_counts.append(counts)
        document_frequency.update(counts)

    total_documents = max(len(chunks), 1)
    indexed: list[KnowledgeChunk] = []
    for chunk, counts in zip(chunks, raw_counts, strict=True):
        weights = {
            token: count * idf(token, document_frequency[token], total_documents)
            for token, count in counts.items()
        }
        indexed.append(
            replace(
                chunk,
                token_counts={token: round(weight, 6) for token, weight in weights.items()},
                norm=round(vector_norm(weights), 6),
            )
        )
    return indexed


def tokenize(text: str) -> list[str]:
    return [match.group(0).lower() for match in TOKEN_RE.finditer(text)]


def idf(token: str, document_frequency: int, total_documents: int) -> float:
    del token
    return math.log((1 + total_documents) / (1 + document_frequency)) + 1.0


def vector_norm(vector: dict[str, float] | Counter[str]) -> float:
    return math.sqrt(sum(float(weight) * float(weight) for weight in vector.values()))


def cosine_score(
    query_vector: dict[str, float],
    query_norm: float,
    chunk_vector: dict[str, float],
    chunk_norm: float,
) -> float:
    if chunk_norm == 0:
        return 0.0
    dot = sum(
        query_vector.get(token, 0.0) * float(chunk_vector.get(token, 0)) for token in query_vector
    )
    return dot / (query_norm * chunk_norm)


def _string_list_from_json(value: object) -> list[str]:
    if isinstance(value, (list, tuple)):
        return [str(item) for item in value]
    return []


def _int_from_json(value: object) -> int:
    if isinstance(value, (str, bytes, bytearray, SupportsInt, SupportsIndex)):
        return int(value)
    raise TypeError(f"Expected JSON scalar convertible to int, got {type(value).__name__}")


def _float_from_json(value: object) -> float:
    if isinstance(value, (str, bytes, bytearray, SupportsFloat, SupportsIndex)):
        return float(value)
    raise TypeError(f"Expected JSON scalar convertible to float, got {type(value).__name__}")


def _float_mapping_from_json(value: object) -> dict[str, float]:
    if not isinstance(value, Mapping):
        return {}
    return {str(key): _float_from_json(item) for key, item in value.items()}


def internal_source_id(project_root: Path, path: Path) -> str:
    try:
        relative = path.relative_to(project_root)
    except ValueError:
        relative = path
    return re.sub(r"[^a-zA-Z0-9_.-]+", "_", str(relative.with_suffix("")))


def internal_topic(path: Path) -> str:
    lowered = str(path).lower()
    if "u3" in lowered:
        return "u3_forecast"
    if "dd" in lowered or "double" in lowered or "laminate" in lowered:
        return "double_double_laminate"
    if "injection" in lowered:
        return "injection_molding"
    return "composites"


def internal_tags(path: Path) -> list[str]:
    tags = ["internal"]
    lowered = str(path).lower()
    if "ppt" in lowered or path.suffix.lower() == ".pptx":
        tags.append("presentation")
    if "report" in lowered:
        tags.append("model-report")
    if "xai" in lowered:
        tags.append("xai")
    if "u3" in lowered:
        tags.append("u3")
    if "dd" in lowered or "double" in lowered:
        tags.append("double-double")
    return tags
